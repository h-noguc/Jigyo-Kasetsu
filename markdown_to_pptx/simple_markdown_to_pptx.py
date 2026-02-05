#!/usr/bin/env python3
"""
シンプルなMarkdownファイルを編集可能なPowerPointファイル（.pptx）に変換するスクリプト

Marpの特殊構文を使わず、標準的なMarkdown形式のみをサポートします。
LLMが生成しやすい形式で、後でPPTXで編集可能なスライドを作成します。

使い方:
    python simple_markdown_to_pptx.py input.md output.pptx

Markdown形式:
    # スライドタイトル
    
    ## 見出し
    
    - リスト項目1
    - リスト項目2
    
    ---
    
    # 次のスライド
"""

import re
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# スタイル設定（一括管理）
STYLE = {
    'background_color': RGBColor(255, 255, 255),   # 白背景
    'text_color': RGBColor(45, 45, 45),            # #2d2d2d - ダークグレー（元の背景色）
    'heading_color': RGBColor(45, 45, 45),         # #2d2d2d - ダークグレー
    'accent_color': RGBColor(45, 45, 45),          # #2d2d2d - ダークグレー
    'quote_color': RGBColor(45, 45, 45),           # #2d2d2d - ダークグレー
    'table_header_bg': RGBColor(240, 240, 240),   # ライトグレー
    'table_cell_bg': RGBColor(255, 255, 255),      # 白
    'font_name': 'TT Commons Pro',                 # メインフォント
    'fallback_font': 'Arial',                     # フォールバックフォント（TT Commons Proが利用できない場合）
    'slide_width': 10,                             # インチ
    'slide_height': 5.625,                         # インチ (16:9 アスペクト比)
}


def parse_simple_markdown(file_path):
    """
    シンプルなMarkdownファイルを解析してスライドのリストを返す
    
    - `---`でスライドを分割
    - YAMLフロントマターはオプション（あっても無視）
    - `<!-- _class: title -->`でタイトルスライドを判定
    - `<!-- _single_page -->`で1ページにまとめる範囲を指定
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # YAMLフロントマターを除去（オプション）
    yaml_match = re.match(r'^---\n(.*?)\n---\n', content, re.DOTALL)
    if yaml_match:
        # YAMLフロントマターの内容を確認（スタイル情報があれば抽出可能）
        content = content[yaml_match.end():]
    
    # スライドを分割（`---`で区切る）
    slide_blocks = re.split(r'\n---\n', content)
    
    # single_pageマークを処理：連続するブロックを1つにまとめる
    slides = []
    i = 0
    while i < len(slide_blocks):
        block = slide_blocks[i].strip()
        if not block:
            i += 1
            continue
        
        # single_pageマークがあるかチェック
        is_single_page = ('<!-- _single_page -->' in block or 
                         '<!--_single_page-->' in block or
                         '<!-- _no_split -->' in block or
                         '<!--_no_split-->' in block)
        
        # タイトルスライドかどうかを判定（Marpの`<!-- _class: title -->`に対応）
        is_title = ('<!-- _class: title -->' in block or 
                   '<!--_class: title-->' in block)
        
        if is_single_page:
            # single_pageマークがある場合、次の`---`までまとめる
            combined_content = [block]
            i += 1
            while i < len(slide_blocks):
                next_block = slide_blocks[i].strip()
                if not next_block:
                    i += 1
                    continue
                # 次のsingle_pageマークやタイトルマークがある場合は終了
                if ('<!-- _single_page -->' in next_block or 
                    '<!--_single_page-->' in next_block or
                    '<!-- _no_split -->' in next_block or
                    '<!--_no_split-->' in next_block or
                    '<!-- _class: title -->' in next_block or
                    '<!--_class: title-->' in next_block):
                    break
                combined_content.append(next_block)
                i += 1
            
            # HTMLコメントを除去
            combined_block = '\n\n'.join(combined_content)
            combined_block = re.sub(r'<!--.*?-->', '', combined_block, flags=re.DOTALL)
            combined_block = combined_block.strip()
            
            if combined_block:
                slides.append({
                    'content': combined_block,
                    'is_title': is_title,
                    'single_page': True
                })
        else:
            # 通常のブロック
            # HTMLコメントを除去（タイトル判定後）
            block = re.sub(r'<!--.*?-->', '', block, flags=re.DOTALL)
            block = block.strip()
            
            if block:
                slides.append({
                    'content': block,
                    'is_title': is_title,
                    'single_page': False
                })
            i += 1
    
    return slides


def parse_markdown_elements(content):
    """
    Markdownコンテンツを解析して構造化データに変換
    
    サポートする要素:
    - 見出し (#, ##, ###)
    - リスト (-, *, 1.)
    - テーブル (| col1 | col2 |)
    - 引用 (>)
    - 通常テキスト
    - 強調 (**bold**, *italic*)
    """
    lines = content.split('\n')
    elements = []
    i = 0
    
    while i < len(lines):
        line = lines[i].strip()
        
        if not line:
            i += 1
            continue
        
        # 見出し
        if line.startswith('#'):
            level = len(line) - len(line.lstrip('#'))
            text = line.lstrip('#').strip()
            elements.append({'type': 'heading', 'level': level, 'text': text})
            i += 1
        
        # 引用（blockquote）
        elif line.startswith('>'):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith('>'):
                quote_text = lines[i].strip().lstrip('>').strip()
                quote_lines.append(quote_text)
                i += 1
            elements.append({'type': 'quote', 'text': '\n'.join(quote_lines)})
            continue
        
        # テーブル
        elif '|' in line:
            table_rows = []
            # ヘッダー行
            header = [cell.strip() for cell in line.split('|')[1:-1]]
            if header:
                table_rows.append(header)
                i += 1
                
                # セパレーター行をスキップ
                if i < len(lines) and re.match(r'^\|[\s\-:]+', lines[i]):
                    i += 1
                
                # データ行
                while i < len(lines) and '|' in lines[i]:
                    row = [cell.strip() for cell in lines[i].split('|')[1:-1]]
                    if row:
                        table_rows.append(row)
                    i += 1
                
                if len(table_rows) > 1:
                    elements.append({'type': 'table', 'data': table_rows})
                continue
        
        # リスト（- または *）
        elif line.startswith('- ') or line.startswith('* '):
            items = []
            base_indent = len(line) - len(line.lstrip())
            while i < len(lines):
                current_line = lines[i].rstrip()
                if not current_line:
                    break
                current_indent = len(current_line) - len(current_line.lstrip())
                if current_line.strip().startswith(('- ', '* ')) and current_indent >= base_indent:
                    item_text = current_line.strip().lstrip('-*').strip()
                    items.append(item_text)
                    i += 1
                else:
                    break
            if items:
                elements.append({'type': 'list', 'items': items})
            continue
        
        # 番号付きリスト
        elif re.match(r'^\d+\.\s', line):
            items = []
            while i < len(lines) and re.match(r'^\d+\.\s', lines[i].strip()):
                item_text = re.sub(r'^\d+\.\s+', '', lines[i].strip())
                items.append(item_text)
                i += 1
            if items:
                elements.append({'type': 'numbered_list', 'items': items})
            continue
        
        # コードブロック（スキップまたはテキストとして扱う）
        elif line.startswith('```'):
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                i += 1
            i += 1
            continue
        
        # 通常のテキスト
        else:
            elements.append({'type': 'text', 'text': line})
            i += 1
    
    return elements


def clean_markdown_text(text):
    """Markdown記法をクリーンアップ（HTMLタグ、<br>など）"""
    # <br>を改行に変換
    text = text.replace('<br>', '\n').replace('<br/>', '\n').replace('<br />', '\n')
    # HTMLタグを削除
    text = re.sub(r'<[^>]+>', '', text)
    return text.strip()


def add_formatted_text(slide, text, left, top, width, height, font_size=18,
                       font_color=STYLE['text_color'], align=PP_ALIGN.LEFT,
                       bold=False, italic=False, auto_resize=True, line_spacing=None):
    """
    Markdown記法（**bold**, *italic*）を含むテキストボックスを追加
    
    auto_resize=Trueの場合、テキストの実際の高さに応じてテキストボックスをリサイズ
    """
    text = clean_markdown_text(text)
    
    # 初期の高さでテキストボックスを作成
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    # 改行で分割
    lines = text.split('\n')
    
    for line_idx, line in enumerate(lines):
        if line_idx > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.alignment = align
        
        # 行間を設定（画像のタイポグラフィガイドに基づく）
        if line_spacing is not None:
            p.line_spacing = line_spacing
        
        # **bold**と*italic*を処理
        # まず**bold**を処理（優先度が高い）
        parts = re.split(r'(\*\*.*?\*\*)', line)
        
        for part in parts:
            if not part:
                continue
            
            run = p.add_run()
            run.font.name = STYLE['font_name']
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            
            if part.startswith('**') and part.endswith('**'):
                # 太字
                run.text = part[2:-2]
                run.font.bold = True
                run.font.italic = italic
            elif part.startswith('*') and part.endswith('*') and len(part) > 2:
                # イタリック（**で囲まれていない場合）
                run.text = part[1:-1]
                run.font.italic = True
                run.font.bold = bold
            else:
                # 通常テキスト（残りの*を削除）
                clean_part = re.sub(r'\*+', '', part)
                run.text = clean_part
                run.font.bold = bold
                run.font.italic = italic
    
    # 自動リサイズ：テキストが収まらない場合は高さを調整
    if auto_resize:
        # テキストフレームの実際の高さを取得
        text_frame.auto_size = None  # 自動サイズを無効化して手動で調整
        
        # 行間を考慮して高さを計算（line_spacingが指定されている場合はそれを使用）
        if line_spacing:
            # line_spacingはPtオブジェクトなので、.pt属性でポイント値を取得
            line_height_pt = line_spacing.pt if hasattr(line_spacing, 'pt') else font_size * 1.5
        else:
            line_height_pt = font_size * 1.5  # デフォルトの行間比
        
        # テキストの実際の行数を計算（折り返しを考慮）
        # テキストボックスの幅（インチ）をポイントに変換
        width_pt = width * 72.0
        
        # 各段落について、折り返し後の行数を計算
        total_lines = 0
        for line in lines:
            if not line.strip():
                # 空行は1行としてカウント
                total_lines += 1
            else:
                # テキストの幅を推定（日本語文字は全角でフォントサイズとほぼ同じ幅）
                # 簡易的な計算：文字数 × フォントサイズ × 0.9（日本語文字の平均幅比）
                char_count = len(line)
                estimated_text_width_pt = char_count * font_size * 0.9
                
                # 折り返し後の行数を計算
                wrapped_lines = max(1, int(estimated_text_width_pt / width_pt) + 1)
                total_lines += wrapped_lines
        
        if total_lines == 0:
            total_lines = 1
        
        # 高さを計算：行数 × 行間 + 余白
        estimated_height = (total_lines * line_height_pt / 72.0) + 0.05  # インチ単位、余白0.05インチ
        
        # テキストボックスの高さを調整（最小でも初期高さの0.5倍、最大で3倍まで）
        min_height = height * 0.5
        max_height = height * 3
        final_height = max(min_height, min(estimated_height, max_height))
        
        text_box.height = Inches(final_height)
        current_height = final_height
    else:
        current_height = height
    
    return text_box, current_height


def add_table(slide, data, left, top, width, height, scale_factor=1.0):
    """テーブルを追加（Markdown記法を処理、scale_factorでフォントサイズを調整）"""
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if rows == 0 or cols == 0:
        return None
    
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            
            # テキストフレームの設定
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.margin_left = Inches(0.1)
            cell.text_frame.margin_right = Inches(0.1)
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)
            
            # セルの内容をクリア
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            
            # 行間を設定 (Sub Text: 11pt, LH: 16pt)
            p.line_spacing = Pt(int(16 * scale_factor))
            
            # Markdown記法を処理
            cell_text = clean_markdown_text(cell_text)
            
            # **bold**と*italic*を処理
            # まず**bold**を処理（優先度が高い）
            parts = re.split(r'(\*\*.*?\*\*)', cell_text)
            
            for part in parts:
                if not part:
                    continue
                
                run = p.add_run()
                run.font.name = STYLE['font_name']
                
                if i == 0:  # ヘッダー行
                    run.font.color.rgb = STYLE['heading_color']
                    run.font.size = Pt(int(11 * scale_factor))
                    run.font.bold = True  # デフォルトで太字
                else:  # データ行
                    run.font.color.rgb = STYLE['text_color']
                    run.font.size = Pt(int(11 * scale_factor))
                    run.font.bold = False
                
                if part.startswith('**') and part.endswith('**'):
                    # 太字
                    run.text = part[2:-2]
                    run.font.bold = True
                    run.font.italic = False
                elif part.startswith('*') and part.endswith('*') and len(part) > 2:
                    # イタリック（**で囲まれていない場合）
                    run.text = part[1:-1]
                    run.font.italic = True
                    run.font.bold = False
                else:
                    # 通常テキスト（残りの*を削除）
                    clean_part = re.sub(r'\*+', '', part)
                    run.text = clean_part
                    run.font.italic = False
                    # ヘッダー行の場合は既に太字、データ行の場合は通常
            
            # 背景色設定
            if i == 0:  # ヘッダー行
                cell.fill.solid()
                cell.fill.fore_color.rgb = STYLE['table_header_bg']
            else:  # データ行
                cell.fill.solid()
                cell.fill.fore_color.rgb = STYLE['table_cell_bg']
    
    return table


def create_slide(prs, slide_data):
    """
    スライドを作成
    
    ルール:
    - is_title=True または H1が最初にある場合 → タイトルスライドとして中央揃え
    - single_page=True の場合 → フォントサイズを調整して1ページに収める
    - それ以外 → 通常スライド
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
    
    # 背景色を設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = STYLE['background_color']
    
    # slide_dataが辞書の場合はcontentとis_titleを取得、文字列の場合は後方互換性のため
    if isinstance(slide_data, dict):
        slide_content = slide_data['content']
        is_title_slide = slide_data.get('is_title', False)
        single_page = slide_data.get('single_page', False)
    else:
        slide_content = slide_data
        is_title_slide = False
        single_page = False
    
    elements = parse_markdown_elements(slide_content)
    
    if not elements:
        return slide
    
    # タイトルスライドの判定（is_titleフラグまたは最初の要素がH1）
    if not is_title_slide:
        is_title_slide = (elements[0]['type'] == 'heading' and elements[0]['level'] == 1)
    
    # スライドの利用可能な高さを計算（マージンを考慮）
    slide_height = STYLE['slide_height']
    if is_title_slide:
        top_margin = 0.5  # タイトルスライドは上部マージン多め
        bottom_margin = 0.5
        available_height = slide_height - top_margin - bottom_margin
        y_pos = 1.0  # タイトルスライドの開始位置を下げる
    else:
        top_margin = 0.5
        bottom_margin = 0.5
        available_height = slide_height - top_margin - bottom_margin
        y_pos = 0.5
    
    # single_pageの場合、まず全体の高さを計算してフォントサイズを調整
    if single_page:
        # 最初のパス：必要な高さを計算
        estimated_height = 0
        for elem in elements:
            if elem['type'] == 'heading':
                level = elem['level']
                if is_title_slide and level == 1:
                    # タイトルスライドのH1は大きめ
                    estimated_height += 1.5
                elif level == 1:
                    estimated_height += 1.2
                elif level == 2:
                    estimated_height += 0.8
                elif level == 3:
                    estimated_height += 0.6
            elif elem['type'] == 'quote':
                estimated_height += 0.8
            elif elem['type'] == 'table':
                estimated_height += min(len(elem['data']) * 0.35, 3.5) + 0.2
            elif elem['type'] == 'list' or elem['type'] == 'numbered_list':
                estimated_height += len(elem['items']) * 0.4 + 0.15
            elif elem['type'] == 'text':
                # タイトルスライドのテキストは小さめに
                if is_title_slide:
                    estimated_height += 0.5
                else:
                    estimated_height += 0.5
        
        # 高さが足りない場合はフォントサイズを調整
        scale_factor = 1.0
        if estimated_height > available_height:
            scale_factor = available_height / estimated_height
            # 最小フォントサイズを設定（10pt以下にはしない）
            min_font_scale = 10.0 / 13.0  # 基準フォント13ptに対して
            scale_factor = max(scale_factor, min_font_scale)
        
        # タイトルスライドの場合はより積極的にスケーリング
        if is_title_slide and scale_factor < 0.9:
            scale_factor = scale_factor * 0.95  # さらに5%縮小
    else:
        scale_factor = 1.0
    
    # 連続するテキストとリストを結合する前処理
    merged_elements = []
    i = 0
    while i < len(elements):
        elem = elements[i]
        
        # テキスト要素の後にリストが続く場合、結合する
        if elem['type'] == 'text' and i + 1 < len(elements):
            next_elem = elements[i + 1]
            if next_elem['type'] == 'list' or next_elem['type'] == 'numbered_list':
                # テキストとリストを結合（視覚的な区切りのため空行を追加）
                if next_elem['type'] == 'list':
                    list_text = '\n'.join(['• ' + item for item in next_elem['items']])
                else:
                    list_text = '\n'.join([f'{idx}. {item}' for idx, item in enumerate(next_elem['items'], 1)])
                
                merged_elem = {
                    'type': 'text_with_list',
                    'text': elem['text'],
                    'list_type': next_elem['type'],
                    'list_items': next_elem['items'],
                    'combined_text': elem['text'] + '\n\n' + list_text  # 空行を追加して視覚的な区切りを作る
                }
                merged_elements.append(merged_elem)
                i += 2  # テキストとリストの両方をスキップ
                continue
        
        merged_elements.append(elem)
        i += 1
    
    elements = merged_elements
    
    for idx, elem in enumerate(elements):
        if elem['type'] == 'heading':
            level = elem['level']
            if is_title_slide and level == 1:
                # タイトルスライドのメインタイトル (H1: 23pt, LH: 36)
                font_size = int(23 * scale_factor)
                line_spacing = Pt(int(36 * scale_factor))  # Line Height: 36pt
                _, actual_height = add_formatted_text(slide, elem['text'], 0.5, y_pos, 9, 1.2 * scale_factor,
                                 font_size=font_size, font_color=STYLE['heading_color'],
                                 align=PP_ALIGN.CENTER, bold=True, line_spacing=line_spacing)
                y_pos += max(actual_height, 1.0 * scale_factor) + 0.3 * scale_factor
            elif level == 1:
                # 通常スライドのH1 (H1: 23pt, LH: 36)
                font_size = int(23 * scale_factor)
                line_spacing = Pt(int(36 * scale_factor))  # Line Height: 36pt
                _, actual_height = add_formatted_text(slide, elem['text'], 0.5, y_pos, 9, 1.0,
                                 font_size=font_size, font_color=STYLE['heading_color'],
                                 bold=True, line_spacing=line_spacing)
                y_pos += actual_height * 1.2
            elif level == 2:
                # H2見出し (H2: 21pt, LH: 32) - 固定サイズで統一
                font_size = 21  # scale_factorを適用せず固定
                line_spacing = Pt(32)  # Line Height: 32pt（固定）
                _, actual_height = add_formatted_text(slide, elem['text'], 0.5, y_pos, 9, 0.8,
                                 font_size=font_size, font_color=STYLE['heading_color'],
                                 bold=True, line_spacing=line_spacing)
                # 次の要素がH3見出しの場合はスペーシングを小さくする
                next_elem = elements[idx + 1] if idx + 1 < len(elements) else None
                if next_elem and next_elem.get('type') == 'heading' and next_elem.get('level') == 3:
                    spacing = 0.05  # H3の場合は小さなスペース
                else:
                    spacing = 0.1  # それ以外は通常のスペース
                y_pos += actual_height + spacing
            elif level == 3:
                # H3見出し (H3: 19pt, LH: 28) - 固定サイズで統一
                font_size = 19  # scale_factorを適用せず固定
                line_spacing = Pt(28)  # Line Height: 28pt（固定）
                _, actual_height = add_formatted_text(slide, elem['text'], 0.5, y_pos, 9, 0.6,
                                 font_size=font_size, font_color=STYLE['accent_color'],
                                 bold=True, line_spacing=line_spacing)
                # 次の要素のタイプを確認してスペーシングを調整
                next_elem = elements[idx + 1] if idx + 1 < len(elements) else None
                if next_elem and next_elem.get('type') == 'text':
                    spacing = 0.05  # テキストの場合は小さなスペース
                else:
                    spacing = 0.08  # それ以外は通常のスペース
                y_pos += actual_height + spacing
        
        elif elem['type'] == 'quote':
            # 引用を左側にボーダー付きで表示
            quote_height = 0.8 * scale_factor
            quote_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(quote_height))
            quote_frame = quote_box.text_frame
            quote_frame.word_wrap = True
            p = quote_frame.paragraphs[0]
            
            # **bold**を処理 (Sub Header 2: 15pt, LH: 24)
            quote_text = clean_markdown_text(elem['text'])
            font_size = int(15 * scale_factor)
            p.line_spacing = Pt(int(24 * scale_factor))  # Line Height: 24pt
            parts = re.split(r'(\*\*.*?\*\*)', quote_text)
            for part in parts:
                if not part:
                    continue
                run = p.add_run()
                run.font.name = STYLE['font_name']
                run.font.size = Pt(font_size)
                run.font.color.rgb = STYLE['quote_color']
                run.font.italic = True
                if part.startswith('**') and part.endswith('**'):
                    run.text = part[2:-2]
                    run.font.bold = True
                else:
                    run.text = re.sub(r'\*+', '', part)
            
            p.alignment = PP_ALIGN.LEFT
            y_pos += quote_height * 1.25
        
        elif elem['type'] == 'table':
            table_height = min(len(elem['data']) * 0.4 * scale_factor, 4.0 * scale_factor)
            add_table(slide, elem['data'], 0.5, y_pos, 9, table_height, scale_factor)
            y_pos += table_height + 0.3 * scale_factor
        
        elif elem['type'] == 'list':
            # リストの全項目を1つのテキストボックスにまとめる
            font_size = int(13 * scale_factor)
            line_spacing = Pt(int(20 * scale_factor))  # Body: 13pt, LH: 20pt
            # 全項目を改行で結合
            list_text = '\n'.join(['• ' + item for item in elem['items']])
            # 高さを推定（項目数 × 行間）
            estimated_height = len(elem['items']) * (20 * scale_factor / 72.0) + 0.2
            _, actual_height = add_formatted_text(slide, list_text, 0.7, y_pos, 8.6, estimated_height,
                                 font_size=font_size, font_color=STYLE['text_color'], line_spacing=line_spacing)
            y_pos += actual_height + 0.2 * scale_factor
        
        elif elem['type'] == 'numbered_list':
            # 番号付きリストの全項目を1つのテキストボックスにまとめる
            font_size = int(13 * scale_factor)
            line_spacing = Pt(int(20 * scale_factor))  # Body: 13pt, LH: 20pt
            # 全項目を改行で結合
            list_text = '\n'.join([f'{idx}. {item}' for idx, item in enumerate(elem['items'], 1)])
            # 高さを推定（項目数 × 行間）
            estimated_height = len(elem['items']) * (20 * scale_factor / 72.0) + 0.2
            _, actual_height = add_formatted_text(slide, list_text, 0.7, y_pos, 8.6, estimated_height,
                                 font_size=font_size, font_color=STYLE['text_color'], line_spacing=line_spacing)
            y_pos += actual_height + 0.2 * scale_factor
        
        elif elem['type'] == 'text_with_list':
            # テキストとリストが結合された要素
            combined_text = clean_markdown_text(elem['combined_text'])
            font_size = int(13 * scale_factor)
            line_spacing = Pt(int(20 * scale_factor))  # Body: 13pt, LH: 20pt
            # 高さを推定（テキスト行数 + リスト項目数）
            estimated_lines = combined_text.count('\n') + 1
            estimated_height = estimated_lines * (20 * scale_factor / 72.0) + 0.2
            _, actual_height = add_formatted_text(slide, combined_text, 0.5, y_pos, 9, estimated_height,
                                 font_size=font_size, font_color=STYLE['text_color'], line_spacing=line_spacing)
            y_pos += actual_height + 0.2 * scale_factor
        
        elif elem['type'] == 'text':
            text = clean_markdown_text(elem['text'])
            if text.strip():
                # タイトルスライドの場合は中央揃え、フォントサイズも調整
                if is_title_slide:
                    if '**' in elem['text']:
                        # Sub Header 1: 17pt, LH: 24pt
                        base_font_size = 17
                        line_spacing = Pt(int(24 * scale_factor))
                    else:
                        # Body: 13pt, LH: 20pt
                        base_font_size = 13
                        line_spacing = Pt(int(20 * scale_factor))
                    font_size = int(base_font_size * scale_factor)
                    font_color = STYLE['text_color'] if '**' in elem['text'] else RGBColor(100, 100, 100)  # グレー（白背景用に調整）
                    _, actual_height = add_formatted_text(slide, text, 0.5, y_pos, 9, 0.5 * scale_factor,
                                     font_size=font_size, font_color=font_color,
                                     align=PP_ALIGN.CENTER, line_spacing=line_spacing)
                    y_pos += max(actual_height, 0.4 * scale_factor) + 0.2 * scale_factor
                else:
                    font_size = int(13 * scale_factor)
                    line_spacing = Pt(int(20 * scale_factor))  # Body: 13pt, LH: 20pt
                    _, actual_height = add_formatted_text(slide, text, 0.5, y_pos, 9, 0.5 * scale_factor,
                                     font_size=font_size, font_color=STYLE['text_color'], line_spacing=line_spacing)
                    y_pos += max(actual_height, 0.5 * scale_factor * 1.2)
        
        # ページの高さを超えたら次のスライドに（single_pageの場合は無視）
        if not single_page and y_pos > (slide_height - bottom_margin):
            break
    
    return slide


def main():
    """メイン関数"""
    if len(sys.argv) < 2:
        print("Usage: python simple_markdown_to_pptx.py <input.md> [output.pptx]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    else:
        # 出力ファイル名を自動生成
        base_name = os.path.splitext(input_file)[0]
        output_file = base_name + '.pptx'
    
    if not os.path.exists(input_file):
        print(f"Error: File not found: {input_file}")
        sys.exit(1)
    
    print(f"Reading Markdown file: {input_file}")
    slides = parse_simple_markdown(input_file)
    print(f"Found {len(slides)} slides")
    
    # PowerPointプレゼンテーションを作成
    prs = Presentation()
    prs.slide_width = Inches(STYLE['slide_width'])
    prs.slide_height = Inches(STYLE['slide_height'])
    
    # 各スライドを作成
    for i, slide_content in enumerate(slides):
        print(f"Creating slide {i+1}/{len(slides)}")
        create_slide(prs, slide_content)
    
    # 保存
    print(f"Saving to: {output_file}")
    prs.save(output_file)
    print("Done!")


if __name__ == "__main__":
    main()
